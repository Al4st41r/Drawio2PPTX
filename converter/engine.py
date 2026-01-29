from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Emu, Pt

from .ppt_map import get_shape_type, get_line_dash, get_arrow_type, get_connector_type
from .utils import hex_to_rgb, px_to_emu, HtmlTextParser, set_line_end


class PptxGenerator:
    def __init__(self, output_file):
        self.output_file = output_file
        self.prs = Presentation()
        self.slide = None
        self.id_to_shape = {}  # Map Draw.io ID to PPTX Shape per slide

    def create_slide(self):
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.id_to_shape = {} # Reset map for new slide
        return self.slide

    def add_vertices(self, vertices):
        if not self.slide:
            self.create_slide()
            
        for v in vertices:
            x, y, w, h = v["x"], v["y"], v["width"], v["height"]
            style = v["style"]

            shape_type = get_shape_type(style)

            shape = self.slide.shapes.add_shape(
                shape_type, px_to_emu(x), px_to_emu(
                    y), px_to_emu(w), px_to_emu(h)
            )

            # Apply Styles
            self._apply_shape_style(shape, style)

            # Add Text
            if v["value"]:
                self._apply_text(shape, v["value"], style)

            self.id_to_shape[v["id"]] = shape

    def add_edges(self, edges):
        for e in edges:
            source_id = e["source"]
            target_id = e["target"]

            if source_id in self.id_to_shape and target_id in self.id_to_shape:
                src_shape = self.id_to_shape[source_id]
                tgt_shape = self.id_to_shape[target_id]
                
                # Determine Connector Type
                # If shapes are aligned, prefer STRAIGHT over ELBOW to prevent auto-routing mess
                conn_type = get_connector_type(e["style"])
                
                if conn_type == MSO_CONNECTOR.ELBOW:
                    # Check alignment
                    src_l, src_r = src_shape.left, src_shape.left + src_shape.width
                    src_t, src_b = src_shape.top, src_shape.top + src_shape.height
                    tgt_l, tgt_r = tgt_shape.left, tgt_shape.left + tgt_shape.width
                    tgt_t, tgt_b = tgt_shape.top, tgt_shape.top + tgt_shape.height
                    
                    x_overlap = max(0, min(src_r, tgt_r) - max(src_l, tgt_l))
                    y_overlap = max(0, min(src_b, tgt_b) - max(src_t, tgt_t))
                    
                    # If significant overlap, use STRAIGHT
                    if x_overlap > min(src_shape.width, tgt_shape.width) * 0.5 or \
                       y_overlap > min(src_shape.height, tgt_shape.height) * 0.5:
                        conn_type = MSO_CONNECTOR.STRAIGHT

                connector = self.slide.shapes.add_connector(
                    conn_type, 0, 0, 0, 0
                )

                # Smart connection logic
                self._connect_shapes(connector, src_shape, tgt_shape, e["style"], conn_type)

                # Apply Styles
                self._apply_line_style(connector.line, e["style"])
                
                # Edge Label (if any)
                if e['value']:
                    self._add_edge_label(e, src_shape, tgt_shape)

    def save(self):
        self.prs.save(self.output_file)

    def _apply_shape_style(self, shape, style):
        # Fill
        fill_color = style.get("fillColor")
        if fill_color == "none":
            shape.fill.background()  # No fill
        elif fill_color:
            rgb = hex_to_rgb(fill_color)
            if rgb:
                shape.fill.solid()
                shape.fill.fore_color.rgb = rgb
        else:
            # Default white fill if not specified
            shape.fill.solid()
            shape.fill.fore_color.rgb = hex_to_rgb("#FFFFFF")

        # Line / Stroke
        stroke_color = style.get("strokeColor")
        if stroke_color == "none":
            shape.line.fill.background()
        elif stroke_color:
            rgb = hex_to_rgb(stroke_color)
            if rgb:
                shape.line.color.rgb = rgb
        else:
            # Default black stroke
            shape.line.color.rgb = hex_to_rgb("#000000")

        # Stroke Width
        if "strokeWidth" in style:
            try:
                width_pt = float(style["strokeWidth"])
                shape.line.width = Pt(width_pt)
            except:
                pass

    def _apply_line_style(self, line, style):
        # Color
        stroke_color = style.get("strokeColor")
        if stroke_color and stroke_color != "none":
            rgb = hex_to_rgb(stroke_color)
            if rgb:
                line.color.rgb = rgb
        else:
            line.color.rgb = hex_to_rgb("#000000")  # Default black

        # Width
        if "strokeWidth" in style:
            try:
                width_pt = float(style["strokeWidth"])
                line.width = Pt(width_pt)
            except:
                pass
        else:
            line.width = Pt(1)
            
        # Dash Style
        line.dash_style = get_line_dash(style)
        
        # Arrowheads (Manual XML injection)
        from .ppt_map import get_arrow_type, get_arrow_size
        
        start_fill = style.get('startFill') != '0'
        end_fill = style.get('endFill') != '0'
        
        start_arrow = get_arrow_type(style.get('startArrow', 'none'), fill=start_fill)
        end_arrow = get_arrow_type(style.get('endArrow', 'none'), fill=end_fill)
        
        # Size mapping
        start_w, start_l = get_arrow_size(style.get('startSize', '6'))
        end_w, end_l = get_arrow_size(style.get('endSize', '6'))
        
        set_line_end(line, head_type=start_arrow, tail_type=end_arrow,
                     head_w=start_w, head_l=start_l, tail_w=end_w, tail_l=end_l)

    def _apply_text(self, shape, text_value, style):
        if not text_value:
            return
            
        tf = shape.text_frame
        tf.word_wrap = True
        tf.margin_left = Emu(0)
        tf.margin_right = Emu(0)
        tf.margin_top = Emu(0)
        tf.margin_bottom = Emu(0)
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        
        parser = HtmlTextParser(text_value)
        segments = parser.parse()
        
        for seg in segments:
            run = p.add_run()
            run.text = seg['text']
            fmt = seg['format']
            
            run.font.bold = fmt['bold']
            run.font.italic = fmt['italic']
            run.font.underline = fmt['underline']
            
            # Color priority: Segment tag > Style attribute > Default Black
            color_hex = fmt.get('color') or style.get('fontColor') or '#000000'
            rgb = hex_to_rgb(color_hex)
            if rgb:
                run.font.color.rgb = rgb
            
            size = fmt.get('size') or style.get('fontSize')
            if size:
                try:
                    run.font.size = Pt(float(size))
                except:
                    pass

    def _connect_shapes(self, connector, src_shape, tgt_shape, edge_style, conn_type):
        def get_idx_from_ratio(x, y):
            try:
                xf, yf = float(x), float(y)
                if yf <= 0.1: return 0 # Top
                if xf >= 0.9: return 1 # Right
                if yf >= 0.9: return 2 # Bottom
                if xf <= 0.1: return 3 # Left
            except: pass
            return None

        # 1. Explicit points
        src_idx = get_idx_from_ratio(edge_style.get('exitX'), edge_style.get('exitY'))
        tgt_idx = get_idx_from_ratio(edge_style.get('entryX'), edge_style.get('entryY'))

        if src_idx is None or tgt_idx is None:
            # 2. Alignment / Overlap Logic
            # Check if shapes are "in the same lane"
            
            src_l, src_r = src_shape.left, src_shape.left + src_shape.width
            src_t, src_b = src_shape.top, src_shape.top + src_shape.height
            tgt_l, tgt_r = tgt_shape.left, tgt_shape.left + tgt_shape.width
            tgt_t, tgt_b = tgt_shape.top, tgt_shape.top + tgt_shape.height
            
            x_overlap = max(0, min(src_r, tgt_r) - max(src_l, tgt_l))
            y_overlap = max(0, min(src_b, tgt_b) - max(src_t, tgt_t))
            
            scx = (src_l + src_r) / 2
            scy = (src_t + src_b) / 2
            tcx = (tgt_l + tgt_r) / 2
            tcy = (tgt_t + tgt_b) / 2
            
            dx = tcx - scx
            dy = tcy - scy
            
            # Determine relationship
            if x_overlap > 0:
                # Vertical alignment
                if dy > 0: # Target Below
                    src_idx = 2; tgt_idx = 0
                else: # Target Above
                    src_idx = 0; tgt_idx = 2
            elif y_overlap > 0:
                # Horizontal alignment
                if dx > 0: # Target Right
                    src_idx = 1; tgt_idx = 3
                else: # Target Left
                    src_idx = 3; tgt_idx = 1
            else:
                # No overlap, use delta heuristic (Diagonal)
                # Use strict cardinal
                if abs(dx) > abs(dy):
                    if dx > 0: src_idx = 1; tgt_idx = 3
                    else: src_idx = 3; tgt_idx = 1
                else:
                    if dy > 0: src_idx = 2; tgt_idx = 0
                    else: src_idx = 0; tgt_idx = 2

        # 3. Optimize Connector Type
        # If perfectly aligned, force STRAIGHT to avoid weird elbows
        # (Only if style wasn't explicitly curved)
        if conn_type == MSO_CONNECTOR.ELBOW:
            # Check if we can use straight line
            can_be_straight = False
            if src_idx == 2 and tgt_idx == 0 and abs(src_shape.left + src_shape.width/2 - (tgt_shape.left + tgt_shape.width/2)) < px_to_emu(10):
                can_be_straight = True # Vertically aligned center
            elif src_idx == 1 and tgt_idx == 3 and abs(src_shape.top + src_shape.height/2 - (tgt_shape.top + tgt_shape.height/2)) < px_to_emu(10):
                can_be_straight = True # Horizontally aligned center
            
            if can_be_straight:
                try:
                    # We can't easily change type of existing object in python-pptx wrapper sometimes, 
                    # but we can try setting the prst property if we accessed xml.
                    # Or simpler: The user already passed `conn_type` to `add_connector`.
                    # We should have decided this BEFORE creating the connector.
                    pass 
                except: pass

        # Calculate connection point coordinates directly
        # idx: 0=top, 1=right, 2=bottom, 3=left
        def conn_point(shape, idx):
            cx = shape.left + shape.width // 2
            cy = shape.top + shape.height // 2
            if idx == 0: return (cx, shape.top)
            if idx == 1: return (shape.left + shape.width, cy)
            if idx == 2: return (cx, shape.top + shape.height)
            if idx == 3: return (shape.left, cy)
            return (cx, cy)

        sx, sy = conn_point(src_shape, src_idx)
        ex, ey = conn_point(tgt_shape, tgt_idx)

        # Set xfrm directly to avoid python-pptx flip normalisation
        from lxml import etree
        nsmap_a = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        nsmap_p = 'http://schemas.openxmlformats.org/presentationml/2006/main'

        spPr = connector._element.find(f'{{{nsmap_p}}}spPr')
        if spPr is None:
            spPr = connector._element.find(f'{{{nsmap_a}}}spPr')
        xfrm = spPr.find(f'{{{nsmap_a}}}xfrm')

        # Determine flips: if source coord > target coord in either axis
        flipH = '1' if sx > ex else None
        flipV = '1' if sy > ey else None

        # Clear existing flip attributes
        for attr in ('flipH', 'flipV'):
            if attr in xfrm.attrib:
                del xfrm.attrib[attr]
        if flipH:
            xfrm.set('flipH', '1')
        if flipV:
            xfrm.set('flipV', '1')

        off = xfrm.find(f'{{{nsmap_a}}}off')
        ext = xfrm.find(f'{{{nsmap_a}}}ext')
        off.set('x', str(min(sx, ex)))
        off.set('y', str(min(sy, ey)))
        ext.set('cx', str(abs(ex - sx)))
        ext.set('cy', str(abs(ey - sy)))

        # Note: stCxn/endCxn references omitted intentionally.
        # LibreOffice recalculates connector geometry from these references,
        # which overrides the explicit xfrm and produces incorrect rendering.


    def _add_edge_label(self, edge, src_shape, tgt_shape):
        # Calculate midpoint using shape centres
        src_cx = src_shape.left + src_shape.width / 2
        src_cy = src_shape.top + src_shape.height / 2
        tgt_cx = tgt_shape.left + tgt_shape.width / 2
        tgt_cy = tgt_shape.top + tgt_shape.height / 2

        geo = edge["geometry"]

        # Relative position along edge (0=source, 0.5=mid, 1=target)
        rel_pos = 0.5
        perp_offset = 0
        pixel_offset_x = 0
        pixel_offset_y = 0

        if geo is not None:
            if geo.get('relative') == '1':
                try:
                    # x is relative position: -1=source, 0=mid, 1=target
                    rel_pos = (float(geo.get('x', 0)) + 1) / 2
                except:
                    pass
                try:
                    perp_offset = float(geo.get('y', 0))
                except:
                    pass

            # Explicit pixel offset from <mxPoint as="offset">
            offset = geo.find("mxPoint")
            if offset is not None and offset.get("as") == "offset":
                try:
                    pixel_offset_x = float(offset.get("x", 0))
                    pixel_offset_y = float(offset.get("y", 0))
                except:
                    pass

        mid_x = src_cx + (tgt_cx - src_cx) * rel_pos
        mid_y = src_cy + (tgt_cy - src_cy) * rel_pos

        # Apply perpendicular offset
        import math
        dx = tgt_cx - src_cx
        dy = tgt_cy - src_cy
        length = math.sqrt(dx * dx + dy * dy)
        if length > 0 and perp_offset != 0:
            # Perpendicular direction (rotate 90 degrees)
            mid_x += (-dy / length) * px_to_emu(perp_offset)
            mid_y += (dx / length) * px_to_emu(perp_offset)

        # Apply pixel offsets
        mid_x += px_to_emu(pixel_offset_x)
        mid_y += px_to_emu(pixel_offset_y)

        # Create a text box
        box_w = px_to_emu(80)
        box_h = px_to_emu(40)

        left = mid_x - (box_w / 2)
        top = mid_y - (box_h / 2)

        tb = self.slide.shapes.add_textbox(left, top, box_w, box_h)
        tb.text_frame.word_wrap = False 
        
        # Style the textbox background (optional, Draw.io labels are usually transparent or white)
        # For now, let's keep it transparent but ensure text is visible
        
        self._apply_text(tb, edge["value"], edge["style"])
