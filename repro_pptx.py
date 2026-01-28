from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor

def create_repro_pptx():
    prs = Presentation()
    # Use a blank slide layout
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Helper to convert pixels to inches (approximate, usually 96 DPI)
    def px_to_in(px):
        return Inches(px / 96.0)

    # 1. Create "Lamp doesn't work" (Rounded Rectangle)
    # XML: x=160, y=80, w=120, h=40
    shape1 = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        px_to_in(160), px_to_in(80), px_to_in(120), px_to_in(40)
    )
    shape1.text_frame.text = "Lamp doesn't work"
    # Basic styling
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(255, 255, 255) # White fill
    shape1.line.color.rgb = RGBColor(0, 0, 0) # Black border

    # 2. Create "Lamp plugged in?" (Rhombus/Diamond)
    # XML: x=170, y=170, w=100, h=80
    shape2 = slide.shapes.add_shape(
        MSO_SHAPE.DIAMOND,
        px_to_in(170), px_to_in(170), px_to_in(100), px_to_in(80)
    )
    shape2.text_frame.text = "Lamp plugged in?"
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape2.line.color.rgb = RGBColor(0, 0, 0)

    # 3. Create Connector (Edge)
    # Connection points: 
    # Rectangles usually have 0 (top), 1 (right), 2 (bottom), 3 (left).
    # We want bottom of shape1 to top of shape2.
    
    # Create the connector object (initial pos doesn't matter much if we connect it)
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW, 
        px_to_in(0), px_to_in(0), px_to_in(10), px_to_in(10)
    )
    
    # Connect them!
    # shape1 (Source) -> Bottom (index 2 likely, need to verify or try)
    # shape2 (Target) -> Top (index 0 likely)
    connector.begin_connect(shape1, 2)
    connector.end_connect(shape2, 0)
    
    # Style the connector
    connector.line.color.rgb = RGBColor(0, 0, 0)
    connector.line.width = Pt(1)

    prs.save('repro_output.pptx')
    print("Generated repro_output.pptx")

if __name__ == "__main__":
    create_repro_pptx()
