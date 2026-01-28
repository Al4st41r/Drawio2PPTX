from pptx import Presentation
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.util import Inches
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

def set_line_end(line, head_type='none', tail_type='none'):
    # Ensure ln exists
    ln = line._get_or_add_ln()
    
    # Namespaces
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    
    # Find or create headEnd
    head = ln.find(f"{{{a_ns}}}headEnd")
    if head is None:
        head = parse_xml(f'<a:headEnd {nsdecls("a")} type="{head_type}" w="med" len="med"/>')
        ln.append(head)
    else:
        head.set('type', head_type)

    # Find or create tailEnd
    tail = ln.find(f"{{{a_ns}}}tailEnd")
    if tail is None:
        tail = parse_xml(f'<a:tailEnd {nsdecls("a")} type="{tail_type}" w="med" len="med"/>')
        ln.append(tail)
    else:
        tail.set('type', tail_type)

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])

connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(1), Inches(1), Inches(5), Inches(5))
line = connector.line

# Attempt to set arrowheads
try:
    set_line_end(line, head_type='triangle', tail_type='oval')
    print("Successfully set line ends.")
except Exception as e:
    print(f"Error: {e}")

prs.save('arrow_test.pptx')
print("Created arrow_test.pptx")