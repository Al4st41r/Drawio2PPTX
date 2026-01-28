import os
import sys
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Add parent directory to path to import main if needed, 
# but we will just test the output file artifact here.

def test_pptx_content(pptx_path):
    print(f"Testing {pptx_path}...")
    if not os.path.exists(pptx_path):
        print("FAIL: Output file not found.")
        return False
        
    prs = Presentation(pptx_path)
    slide = prs.slides[0]
    
    shapes = list(slide.shapes)
    print(f"Found {len(shapes)} shapes on slide 1.")
    
    # We expect: 
    # Shapes: 
    # 1. Lamp doesn't work (Rect)
    # 2. Lamp plugged in? (Diamond)
    # 3. Plug in lamp (Rect)
    # 4. Bulb burned out? (Diamond)
    # 5. Repair Lamp (Rect)
    # 6. Replace Bulb (Rect)
    # Total 6 shapes.
    
    # Connectors:
    # 1. Lamp -> Plugged?
    # 2. Plugged? -> Yes -> Bulb burned?
    # 3. Plugged? -> No -> Plug in
    # 4. Bulb burned? -> Yes -> Replace
    # 5. Bulb burned? -> No -> Repair
    # Total 5 connectors.
    
    # TextBoxes (Labels on connectors):
    # "Yes", "No", "Yes", "No" (4 labels)
    
    # Total expected roughly 15 items.
    
    shape_count = 0
    connector_count = 0
    text_found = False
    
    for shape in shapes:
        # Debug print
        # print(f"Shape type: {shape.shape_type}")
        
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            shape_count += 1
            if shape.has_text_frame:
                if "Lamp doesn't work" in shape.text_frame.text:
                    text_found = True
                    print("  - Found text: 'Lamp doesn't work'")
        elif shape.shape_type == 10: # MSO_SHAPE_TYPE.LINE is 9, but Connectors are often treated differently?
             # Let's check for LINE (9) which covers connectors usually
             pass
        
        # Check for connector explicitly
        # In python-pptx, connectors often report as LINE (9)
        if shape.shape_type == MSO_SHAPE_TYPE.LINE:
            connector_count += 1
        elif shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
            # Our labels might be text boxes
            pass
            
    print(f"  - AutoShapes: {shape_count}")
    print(f"  - Connectors: {connector_count}")
    
    if shape_count < 6:
        print("FAIL: Too few shapes found.")
        return False
        
    if connector_count < 5:
        print("FAIL: Too few connectors found.")
        return False
        
    if not text_found:
        print("FAIL: Specific text not found.")
        return False
        
    print("PASS: Basic validation successful.")
    return True

if __name__ == "__main__":
    if len(sys.argv) > 1:
        pptx_path = sys.argv[1]
    else:
        pptx_path = "output_test.pptx"
        
    if not test_pptx_content(pptx_path):
        sys.exit(1)
