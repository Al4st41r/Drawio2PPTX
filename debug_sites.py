from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.util import Inches

def create_debug_pptx():
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Diamond
    shape1 = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, Inches(2), Inches(2), Inches(2), Inches(2))
    shape1.text_frame.text = "DIAMOND"
    
    # 4 small rects to connect to
    r_top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(1), Inches(0.5), Inches(0.5))
    r_right = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3), Inches(0.5), Inches(0.5))
    r_bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(3), Inches(5), Inches(0.5), Inches(0.5))
    r_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(0.5), Inches(0.5))
    
    # Connect to index 0
    c0 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0,0,0,0)
    c0.begin_connect(shape1, 0)
    c0.end_connect(r_top, 0)
    
    # Connect to index 1
    c1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0,0,0,0)
    c1.begin_connect(shape1, 1)
    c1.end_connect(r_right, 0)
    
    # Connect to index 2
    c2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0,0,0,0)
    c2.begin_connect(shape1, 2)
    c2.end_connect(r_bottom, 0)
    
    # Connect to index 3
    c3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, 0,0,0,0)
    c3.begin_connect(shape1, 3)
    c3.end_connect(r_left, 0)
    
    prs.save('debug_sites.pptx')
    print("Created debug_sites.pptx")

if __name__ == "__main__":
    create_debug_pptx()
